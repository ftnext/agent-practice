from __future__ import annotations

import unicodedata
from io import BytesIO
from typing import TYPE_CHECKING, Literal
from zipfile import ZipFile

from pydantic import BaseModel, Field, model_validator

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
    from pptx.slide import Slide


_NAVY = (15, 32, 58)
_BLUE = (38, 116, 217)
_PALE_BLUE = (232, 241, 252)
_WHITE = (255, 255, 255)
_INK = (32, 39, 48)
_MUTED = (91, 101, 115)
_FONT_NAME = "Noto Sans JP"


class SlideSpec(BaseModel):
    """Content and layout for one slide."""

    kind: Literal["title", "section", "content", "two_column", "conclusion"]
    title: str = Field(min_length=1, max_length=120)
    subtitle: str = Field(default="", max_length=240)
    bullets: list[str] = Field(default_factory=list, max_length=7)
    right_title: str = Field(default="", max_length=80)
    right_bullets: list[str] = Field(default_factory=list, max_length=7)
    key_message: str = Field(default="", max_length=240)
    sources: list[str] = Field(default_factory=list, max_length=3)

    @model_validator(mode="after")
    def validate_slide_title_width(self) -> SlideSpec:
        """Keep non-title headings on one line in the fixed slide layout."""
        display_width = sum(
            2 if unicodedata.east_asian_width(character) in "FWA" else 1
            for character in self.title
        )
        if self.kind != "title" and display_width > 46:
            raise ValueError(
                "Non-title slide headings must be at most 46 display columns."
            )
        return self


class DeckSpec(BaseModel):
    """A researched presentation ready for deterministic rendering."""

    title: str = Field(min_length=1, max_length=120)
    subtitle: str = Field(default="", max_length=240)
    slides: list[SlideSpec] = Field(min_length=2, max_length=20)


def _set_background(slide: Slide, color: tuple[int, int, int]) -> None:
    """Set a solid slide background."""
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _add_text_box(
    slide: Slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    alignment: str = "left",
) -> None:
    """Add consistently styled text to a slide."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = _FONT_NAME
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)


def _add_bullets(
    slide: Slide,
    bullets: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a compact bullet list to a slide."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = (
            text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        )
        paragraph.text = f"• {bullet}"
        paragraph.font.name = _FONT_NAME
        paragraph.font.size = Pt(21)
        paragraph.font.color.rgb = RGBColor(*_INK)
        paragraph.space_after = Pt(11)


def _add_accent(slide: Slide) -> None:
    """Add the common top accent bar."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.11),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*_BLUE)
    accent.line.fill.background()


def _add_sources(
    slide: Slide,
    sources: list[str],
    *,
    color: tuple[int, int, int] = _MUTED,
) -> None:
    """Add source URLs to the footer and speaker notes."""
    if not sources:
        return
    source_text = "Sources: " + " | ".join(sources)
    _add_text_box(
        slide,
        source_text,
        left=0.7,
        top=7.05,
        width=11.9,
        height=0.28,
        font_size=7,
        color=color,
    )
    slide.notes_slide.notes_text_frame.text = "[Sources]\n" + "\n".join(
        f"- {source}" for source in sources
    )


def _render_title_slide(presentation: PresentationType, spec: SlideSpec) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, _NAVY)
    _add_text_box(
        slide,
        spec.title,
        left=0.85,
        top=1.55,
        width=11.2,
        height=1.8,
        font_size=50,
        color=_WHITE,
        bold=True,
    )
    if spec.subtitle:
        _add_text_box(
            slide,
            spec.subtitle,
            left=0.9,
            top=3.55,
            width=10.8,
            height=0.9,
            font_size=24,
            color=_PALE_BLUE,
        )
    _add_text_box(
        slide,
        "RESEARCHED WITH GOOGLE SEARCH",
        left=0.9,
        top=6.45,
        width=5.4,
        height=0.35,
        font_size=10,
        color=_PALE_BLUE,
    )
    _add_sources(slide, spec.sources, color=_PALE_BLUE)


def _render_section_slide(presentation: PresentationType, spec: SlideSpec) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, _BLUE)
    _add_text_box(
        slide,
        spec.title,
        left=1.0,
        top=2.3,
        width=11.3,
        height=1.2,
        font_size=40,
        color=_WHITE,
        bold=True,
        alignment="center",
    )
    if spec.subtitle:
        _add_text_box(
            slide,
            spec.subtitle,
            left=1.5,
            top=3.65,
            width=10.3,
            height=0.8,
            font_size=24,
            color=_WHITE,
            alignment="center",
        )
    _add_sources(slide, spec.sources, color=_PALE_BLUE)


def _render_content_slide(presentation: PresentationType, spec: SlideSpec) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, _WHITE)
    _add_accent(slide)
    _add_text_box(
        slide,
        spec.title,
        left=0.7,
        top=0.48,
        width=11.9,
        height=0.85,
        font_size=35,
        color=_NAVY,
        bold=True,
    )
    _add_bullets(
        slide,
        spec.bullets,
        left=0.85,
        top=1.55,
        width=11.45,
        height=4.7,
    )
    if spec.key_message:
        _add_text_box(
            slide,
            spec.key_message,
            left=0.9,
            top=6.15,
            width=11.3,
            height=0.58,
            font_size=17,
            color=_BLUE,
            bold=True,
            alignment="center",
        )
    _add_sources(slide, spec.sources)


def _render_two_column_slide(presentation: PresentationType, spec: SlideSpec) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, _WHITE)
    _add_accent(slide)
    _add_text_box(
        slide,
        spec.title,
        left=0.7,
        top=0.48,
        width=11.9,
        height=0.85,
        font_size=35,
        color=_NAVY,
        bold=True,
    )
    for left in (0.7, 6.75):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(1.5),
            Inches(5.85),
            Inches(4.95),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_PALE_BLUE)
        card.line.color.rgb = RGBColor(198, 218, 243)
    _add_text_box(
        slide,
        spec.subtitle or "ポイント",
        left=1.0,
        top=1.78,
        width=5.2,
        height=0.5,
        font_size=24,
        color=_NAVY,
        bold=True,
    )
    _add_bullets(
        slide,
        spec.bullets,
        left=1.0,
        top=2.45,
        width=5.1,
        height=3.55,
    )
    _add_text_box(
        slide,
        spec.right_title or "比較",
        left=7.05,
        top=1.78,
        width=5.1,
        height=0.5,
        font_size=24,
        color=_NAVY,
        bold=True,
    )
    _add_bullets(
        slide,
        spec.right_bullets,
        left=7.05,
        top=2.45,
        width=5.0,
        height=3.55,
    )
    _add_sources(slide, spec.sources)


def _render_conclusion_slide(presentation: PresentationType, spec: SlideSpec) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide, _NAVY)
    _add_text_box(
        slide,
        spec.title,
        left=0.9,
        top=0.75,
        width=11.5,
        height=0.9,
        font_size=35,
        color=_WHITE,
        bold=True,
        alignment="center",
    )
    for index, bullet in enumerate(spec.bullets[:3]):
        _add_text_box(
            slide,
            bullet,
            left=1.25,
            top=2.0 + index * 1.15,
            width=10.8,
            height=0.75,
            font_size=24,
            color=_WHITE,
            bold=index == 0,
            alignment="center",
        )
    if spec.key_message:
        _add_text_box(
            slide,
            spec.key_message,
            left=1.0,
            top=5.85,
            width=11.3,
            height=0.7,
            font_size=19,
            color=_PALE_BLUE,
            bold=True,
            alignment="center",
        )
    _add_sources(slide, spec.sources, color=_PALE_BLUE)


def render_presentation(spec: DeckSpec) -> bytes:
    """Render and verify a presentation in a worker thread."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    renderers = {
        "title": _render_title_slide,
        "section": _render_section_slide,
        "content": _render_content_slide,
        "two_column": _render_two_column_slide,
        "conclusion": _render_conclusion_slide,
    }
    for slide_spec in spec.slides:
        renderers[slide_spec.kind](presentation, slide_spec)

    output = BytesIO()
    presentation.save(output)
    payload = output.getvalue()

    with ZipFile(BytesIO(payload)) as archive:
        broken_entry = archive.testzip()
        if broken_entry is not None:
            raise ValueError(f"Generated PPTX contains a broken entry: {broken_entry}")

    verified = Presentation(BytesIO(payload))
    if len(verified.slides) != len(spec.slides):
        raise ValueError("Generated PPTX has an unexpected slide count.")
    return payload
